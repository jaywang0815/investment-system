"""
PPT export — generates a presentation with one chart slide per ticker.
"""
import io
import math
import unicodedata
from datetime import date
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from matplotlib.patches import Rectangle as MplRect

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── colours ──────────────────────────────────────────────────────
_NAVY    = RGBColor(0x0F, 0x25, 0x57)
_BLUE    = RGBColor(0x1E, 0x3A, 0x8A)
_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
_GRAY    = RGBColor(0x64, 0x74, 0x8B)
_GREEN   = RGBColor(0x16, 0xA3, 0x4A)
_RED     = RGBColor(0xDC, 0x26, 0x26)
# reference line colours — maximally distinct
_C_INIT  = RGBColor(0xFF, 0xD7, 0x00)   # gold/yellow  — 期初
_C_KO    = RGBColor(0x00, 0xE6, 0x76)   # bright green — KO
_C_KI    = RGBColor(0xFF, 0x33, 0x33)   # bright red   — KI


def _clean_ticker(t: str) -> str:
    """Normalize full-width chars (ＡＮＥＴ→ANET), strip $, uppercase."""
    return unicodedata.normalize("NFKC", t).lstrip("$").strip().upper()


def _is_valid(val) -> bool:
    """Return True only if val is a real non-NaN number."""
    if val is None:
        return False
    try:
        return not math.isnan(float(val))
    except (TypeError, ValueError):
        return bool(val)


def _textbox(slide, left, top, width, height, text, size,
             bold=False, color=_WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _rect(slide, left, top, width, height, fill_rgb):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape


# ── indicator helpers ─────────────────────────────────────────────
def _calc_rsi(close, period=14):
    delta = close.diff()
    gain = delta.clip(lower=0).rolling(period).mean()
    loss = (-delta.clip(upper=0)).rolling(period).mean()
    return 100 - (100 / (1 + gain / loss))


def _calc_macd(close, fast=12, slow=26, signal=9):
    macd_line   = close.ewm(span=fast, adjust=False).mean() - \
                  close.ewm(span=slow, adjust=False).mean()
    signal_line = macd_line.ewm(span=signal, adjust=False).mean()
    hist        = macd_line - signal_line
    return macd_line, signal_line, hist


def _strip_tz(idx):
    """Remove timezone from DatetimeIndex safely."""
    if hasattr(idx, "tz") and idx.tz is not None:
        try:
            return idx.tz_localize(None)
        except Exception:
            return idx.tz_convert(None)
    return idx


# ── mplfinance chart ──────────────────────────────────────────────
def _chart_mplfinance(ohlcv, ticker, hlines: dict | None = None) -> bytes | None:
    try:
        import mplfinance as mpf

        close = ohlcv["Close"].astype(float)
        rsi = _calc_rsi(close).fillna(50)
        macd_line, signal_line, macd_hist = _calc_macd(close)
        macd_line   = macd_line.fillna(0)
        signal_line = signal_line.fillna(0)
        macd_hist   = macd_hist.fillna(0)
        hist_colors = ["#26a69a" if v >= 0 else "#ef5350"
                       for v in macd_hist]

        last  = close.iloc[-1]
        prev  = close.iloc[-2] if len(close) > 1 else last
        chg   = last - prev
        chg_p = chg / prev * 100
        sign  = "+" if chg >= 0 else ""
        title_str = (f"\n{ticker}   ${last:,.2f}   "
                     f"{sign}{chg:.2f} ({sign}{chg_p:.2f}%)")

        mc = mpf.make_marketcolors(
            up="#26a69a", down="#ef5350",
            wick={"up": "#26a69a", "down": "#ef5350"},
            edge={"up": "#26a69a", "down": "#ef5350"},
            volume={"up": "#26a69a", "down": "#ef5350"},
        )
        style = mpf.make_mpf_style(
            marketcolors=mc,
            gridstyle="--",
            gridcolor="#E2E8F0",
            gridaxis="both",
            facecolor="white",
            figcolor="white",
            rc={
                "axes.labelcolor": "#64748B",
                "xtick.color": "#64748B",
                "ytick.color": "#64748B",
                "font.size": 9,
            },
        )

        apds = [
            mpf.make_addplot([70] * len(rsi), panel=2,
                             color="#cbd5e1", linestyle="--", width=0.7),
            mpf.make_addplot([30] * len(rsi), panel=2,
                             color="#cbd5e1", linestyle="--", width=0.7),
            mpf.make_addplot(rsi, panel=2, color="#9333ea",
                             width=1.4, ylabel="RSI 14"),
            mpf.make_addplot(macd_hist, panel=3, type="bar",
                             color=hist_colors, alpha=0.75, ylabel="MACD"),
            mpf.make_addplot(macd_line,   panel=3, color="#3B82F6", width=1.2),
            mpf.make_addplot(signal_line, panel=3, color="#F97316", width=1.2),
        ]

        fig, axes = mpf.plot(
            ohlcv,
            type="candle",
            style=style,
            volume=True,
            addplot=apds,
            returnfig=True,
            figsize=(14, 9),
            panel_ratios=(4, 1.2, 1.5, 1.5),
            title=title_str,
        )

        axes[0].title.set_color("#1E3A8A")
        axes[0].title.set_fontsize(13)
        axes[0].title.set_fontweight("bold")

        _draw_hlines(axes[0], hlines, len(ohlcv))

        # RSI shaded zones — axes order: [candle, candle_twin, vol, vol_twin, rsi, ...]
        try:
            rsi_ax = axes[4]
            rsi_ax.axhspan(70, 100, alpha=0.07, color="#ef5350", zorder=0)
            rsi_ax.axhspan(0,  30,  alpha=0.07, color="#26a69a", zorder=0)
            rsi_ax.set_ylim(0, 100)
        except Exception:
            pass

        fig.tight_layout(pad=1.0)
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150,
                    bbox_inches="tight", facecolor="white")
        plt.close(fig)
        buf.seek(0)
        return buf.read()
    except Exception as e:
        print(f"[mplfinance error] {e}")
        return None


# ── fallback: pure matplotlib candlestick + RSI + MACD (3-panel) ─
def _chart_matplotlib(ohlcv, ticker, hlines: dict | None = None) -> bytes | None:
    try:
        close  = ohlcv["Close"].astype(float)
        opens  = ohlcv["Open"].astype(float)
        highs  = ohlcv["High"].astype(float)
        lows   = ohlcv["Low"].astype(float)
        dates  = ohlcv.index
        n      = len(dates)
        xs     = list(range(n))

        rsi = _calc_rsi(close).fillna(50)
        macd_line, signal_line, macd_hist = _calc_macd(close)
        macd_line   = macd_line.fillna(0)
        signal_line = signal_line.fillna(0)
        macd_hist   = macd_hist.fillna(0)

        fig, axes = plt.subplots(
            3, 1, figsize=(14, 8),
            gridspec_kw={"height_ratios": [5, 1.8, 1.8]},
            facecolor="white", sharex=True,
        )
        ax1, ax2, ax3 = axes

        for ax in axes:
            ax.set_facecolor("#FAFBFD")
            ax.grid(axis="both", color="#E8EDF2", linewidth=0.6, linestyle="--")
            ax.spines["top"].set_visible(False)
            ax.spines["left"].set_visible(False)
            ax.spines["right"].set_color("#E2E8F0")
            ax.spines["bottom"].set_color("#E2E8F0")
            ax.tick_params(colors="#64748B", labelsize=9)
            ax.yaxis.tick_right()

        # ── Candlestick ──────────────────────────────────────────
        cl_arr = close.values
        op_arr = opens.values
        hi_arr = highs.values
        lo_arr = lows.values
        for i in xs:
            is_up   = cl_arr[i] >= op_arr[i]
            col     = "#26a69a" if is_up else "#ef5350"
            body_lo = min(op_arr[i], cl_arr[i])
            body_hi = max(op_arr[i], cl_arr[i])
            ax1.plot([i, i], [lo_arr[i], hi_arr[i]], color=col, linewidth=0.9, zorder=1)
            ax1.add_patch(MplRect(
                (i - 0.4, body_lo), 0.8, max(body_hi - body_lo, 0.01),
                color=col, zorder=2,
            ))
        ax1.set_xlim(-1, n)
        ax1.set_ylabel("Price", color="#94A3B8", fontsize=9, labelpad=6)
        _draw_hlines(ax1, hlines, n)

        # ── RSI ──────────────────────────────────────────────────
        ax2.plot(xs, rsi.values, color="#9333ea", linewidth=1.4)
        ax2.axhline(70, color="#ef5350", linestyle="--", linewidth=0.8, alpha=0.7)
        ax2.axhline(30, color="#26a69a", linestyle="--", linewidth=0.8, alpha=0.7)
        ax2.axhspan(70, 100, alpha=0.05, color="#ef5350")
        ax2.axhspan(0,  30,  alpha=0.05, color="#26a69a")
        ax2.set_ylim(0, 100)
        ax2.set_yticks([30, 50, 70])
        ax2.text(1, 72, "RSI 14", transform=ax2.get_yaxis_transform(),
                 color="#9333ea", fontsize=8, va="bottom")

        # ── MACD ─────────────────────────────────────────────────
        mh = macd_hist.values
        bar_cols = ["#26a69a" if v >= 0 else "#ef5350" for v in mh]
        ax3.bar(xs, mh, color=bar_cols, alpha=0.6, width=0.8)
        ax3.plot(xs, macd_line.values,   color="#3B82F6", linewidth=1.3)
        ax3.plot(xs, signal_line.values, color="#F97316", linewidth=1.3)
        ax3.axhline(0, color="#CBD5E1", linewidth=0.6)
        ax3.text(1, 0.95, "MACD", transform=ax3.get_yaxis_transform(),
                 color="#3B82F6", fontsize=8, va="top")

        # x-axis date labels
        tick_step = max(1, n // 10)
        tick_pos  = list(range(0, n, tick_step))
        ax3.set_xticks(tick_pos)
        span_days = (dates[-1] - dates[0]).days if n > 1 else 0
        fmt = "%Y/%m" if span_days > 180 else "%m/%d"
        ax3.set_xticklabels(
            [dates[i].strftime(fmt) for i in tick_pos],
            rotation=0, fontsize=9, color="#64748B",
        )

        plt.tight_layout(pad=0.5)
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150,
                    bbox_inches="tight", facecolor="white")
        plt.close(fig)
        buf.seek(0)
        return buf.read()
    except Exception as e:
        print(f"[matplotlib chart error] {ticker}: {type(e).__name__}: {e}")
        return None


# ── last-resort: simple line chart ───────────────────────────────
def _chart_simple(ohlcv, ticker, hlines: dict | None = None) -> bytes | None:
    try:
        close = ohlcv["Close"].astype(float)
        dates = ohlcv.index
        n     = len(dates)
        xs    = list(range(n))

        last  = float(close.iloc[-1])
        prev  = float(close.iloc[-2]) if n > 1 else last
        chg_p = (last - prev) / prev * 100 if prev != 0 else 0
        sign  = "+" if chg_p >= 0 else ""
        color = "#26a69a" if chg_p >= 0 else "#ef5350"

        fig, ax = plt.subplots(figsize=(14, 6), facecolor="white")
        ax.set_facecolor("white")
        ax.plot(xs, close.values, color=color, linewidth=2)
        ax.fill_between(xs, close.values, close.min() * 0.995,
                        alpha=0.1, color=color)
        ax.set_title(
            f"{ticker}   ${last:,.2f}   {sign}{chg_p:.2f}%",
            fontsize=13, color="#1E3A8A", fontweight="bold",
        )
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.grid(axis="y", color="#E2E8F0", linewidth=0.7, linestyle="--")
        ax.tick_params(colors="#64748B", labelsize=8)
        ax.yaxis.tick_right()
        tick_step = max(1, n // 8)
        tick_pos  = list(range(0, n, tick_step))
        ax.set_xticks(tick_pos)
        ax.set_xticklabels(
            [dates[i].strftime("%m/%d") for i in tick_pos],
            rotation=0, fontsize=8, color="#64748B",
        )
        _draw_hlines(ax, hlines, n)
        plt.tight_layout(pad=0.8)
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150,
                    bbox_inches="tight", facecolor="white")
        plt.close(fig)
        buf.seek(0)
        return buf.read()
    except Exception as e:
        print(f"[simple chart error] {ticker}: {e}")
        return None


def _draw_hlines(ax, hlines: dict, n: int) -> None:
    if not hlines:
        return
    init_p = hlines.get("initial")
    ko_p   = hlines.get("ko")
    ki_p   = hlines.get("ki")

    def _hline(price, color, label):
        if not (price and _is_valid(price)):
            return
        ax.axhline(float(price), color=color, linestyle="--",
                   linewidth=2.5, alpha=1.0, zorder=5)
        ax.text(n - 1, float(price),
                f" {label}  ${float(price):,.2f} ",
                va="bottom", ha="right",
                color="#222222", fontsize=10, fontweight="bold", zorder=6,
                bbox=dict(facecolor=color, edgecolor=color,
                          alpha=0.92, boxstyle="round,pad=0.3", linewidth=0))

    ko_eq_init = ko_p and init_p and abs(float(ko_p) - float(init_p)) < 0.01
    ki_eq_init = ki_p and init_p and abs(float(ki_p) - float(init_p)) < 0.01

    if ko_eq_init:
        _hline(init_p, "#00E676", "KO = INIT")
    else:
        _hline(init_p, "#FFD700", "INIT")
        _hline(ko_p,   "#00E676", "KO")

    if not ki_eq_init:
        _hline(ki_p, "#FF3333", "KI")


def _make_chart_png(ticker: str, period: str = "6mo",
                    hlines: dict | None = None) -> bytes | None:
    try:
        import yfinance as yf
        from datetime import datetime, timedelta
        ticker = _clean_ticker(ticker)
        tk = yf.Ticker(ticker)
        if period == "18mo":
            start = (datetime.today() - timedelta(days=548)).strftime("%Y-%m-%d")
            hist = tk.history(start=start)
        else:
            hist = tk.history(period=period)
        if hist.empty or len(hist) < 10:
            print(f"[chart] {ticker}: not enough data ({len(hist)} rows)")
            return None

        hist.index = _strip_tz(hist.index)
        hist = hist.dropna(subset=["Close", "Open", "High", "Low"])
        if len(hist) < 10:
            print(f"[chart] {ticker}: not enough clean data after dropna")
            return None

        ohlcv = hist[["Open", "High", "Low", "Close", "Volume"]].copy()
        ohlcv["Volume"] = ohlcv["Volume"].fillna(0)

        # resample to weekly for long periods so candles are readable
        if period in ("1y", "18mo", "2y", "3y", "4y", "5y"):
            ohlcv = ohlcv.resample("W").agg({
                "Open": "first", "High": "max",
                "Low": "min", "Close": "last", "Volume": "sum",
            }).dropna()

        result = _chart_mplfinance(ohlcv, ticker, hlines=hlines)
        if result is None:
            print(f"[chart] {ticker}: mplfinance failed, trying matplotlib")
            result = _chart_matplotlib(ohlcv, ticker, hlines=hlines)
        if result is None:
            print(f"[chart] {ticker}: matplotlib failed, trying simple line")
            result = _chart_simple(ohlcv, ticker, hlines=hlines)
        if result is None:
            print(f"[chart] {ticker}: all chart methods failed")
        return result
    except Exception as e:
        print(f"[chart error] {ticker}: {type(e).__name__}: {e}")
        return None


def build_ppt(tickers: list[str], sn_info: dict | None = None,
              period: str = "6mo") -> bytes:
    """
    tickers  : list of stock symbols
    sn_info  : dict[ticker] = {ko, ki, product_code, ...}  (optional)
    period   : yfinance period string e.g. "1mo", "3mo", "6mo", "1y", "2y"
    Returns  : PPT file as bytes
    """
    W = Inches(13.33)
    H = Inches(7.5)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]  # fully blank

    # ── Title slide ───────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = _NAVY

    # thin top accent line
    _rect(s, 0, 0, W, Inches(0.06), RGBColor(0x1E, 0x3A, 0x8A))

    # main title — centred vertically
    _textbox(s, Inches(0), Inches(2.4), W, Inches(1.5),
             "投資組合持倉分析報告",
             46, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)

    # english subtitle
    _textbox(s, Inches(0), Inches(3.95), W, Inches(0.6),
             "Investment Portfolio Holdings Analysis Report",
             14, color=RGBColor(0x7B, 0xA7, 0xD9), align=PP_ALIGN.CENTER)

    # thin horizontal rule
    _rect(s, Inches(2.5), Inches(4.7), Inches(8.33), Inches(0.025),
          RGBColor(0x2D, 0x4E, 0x88))

    # bottom info strip
    _rect(s, 0, Inches(6.5), W, Inches(1.0), RGBColor(0x0A, 0x1A, 0x3D))
    _textbox(s, Inches(0.5), Inches(6.62), Inches(6), Inches(0.45),
             date.today().strftime("%Y  /  %m  /  %d"),
             16, color=RGBColor(0x94, 0xA3, 0xB8))
    _textbox(s, Inches(7), Inches(6.62), Inches(6), Inches(0.45),
             f"共  {len(tickers)}  個標的",
             16, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.RIGHT)

    # ── fetch current price (reliable: use history not fast_info) ──
    def _fetch_price(tk_str: str):
        try:
            import yfinance as yf
            hist = yf.Ticker(_clean_ticker(tk_str)).history(period="5d")
            if len(hist) >= 2:
                price = float(hist["Close"].iloc[-1])
                prev  = float(hist["Close"].iloc[-2])
                chg   = price - prev
                chg_p = chg / prev * 100 if prev else 0.0
                return price, chg, chg_p
            elif len(hist) == 1:
                return float(hist["Close"].iloc[-1]), 0.0, 0.0
        except Exception:
            pass
        return None, None, None

    # ── One slide per ticker ──────────────────────────────────────
    for ticker in tickers:
        s = prs.slides.add_slide(blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = _WHITE

        # ── SN info ───────────────────────────────────────────────
        info   = (sn_info or {}).get(ticker, {})
        ko     = info.get("ko")
        ki     = info.get("ki")
        init_p = info.get("initial_price")
        code   = info.get("product_code", "")

        ko_price = round(float(init_p) * float(ko), 2) if _is_valid(init_p) and _is_valid(ko) else None
        ki_price = round(float(init_p) * float(ki), 2) if _is_valid(init_p) and _is_valid(ki) else None

        # ── header bar ────────────────────────────────────────────
        _rect(s, 0, 0, W, Inches(0.85), _BLUE)
        _textbox(s, Inches(0.3), Inches(0.06), Inches(7), Inches(0.72),
                 ticker, 36, bold=True, color=_WHITE)
        _textbox(s, Inches(7.5), Inches(0.17), Inches(5.5), Inches(0.5),
                 date.today().strftime("%Y/%m/%d"),
                 13, color=RGBColor(0xBF, 0xDB, 0xFF), align=PP_ALIGN.RIGHT)

        # ── info row (navy bg) ────────────────────────────────────
        _rect(s, 0, Inches(0.85), W, Inches(1.05), _NAVY)

        curr, chg, chg_p = _fetch_price(ticker)

        # current price + change (left block)
        if curr is not None:
            chg_col = RGBColor(0x4a, 0xde, 0x80) if chg >= 0 else RGBColor(0xf8, 0x71, 0x71)
            arrow   = "▲" if chg >= 0 else "▼"
            sign    = "+" if chg >= 0 else ""
            _textbox(s, Inches(0.3), Inches(0.88), Inches(3.2), Inches(0.6),
                     f"${curr:,.2f}", 26, bold=True, color=_WHITE)
            _textbox(s, Inches(0.3), Inches(1.44), Inches(4.5), Inches(0.4),
                     f"{arrow} {abs(chg):.2f}  ({sign}{chg_p:.2f}%)", 12, color=chg_col)
        else:
            _textbox(s, Inches(0.3), Inches(0.88), Inches(3.2), Inches(0.6),
                     "—", 24, color=RGBColor(0x64, 0x74, 0x8B))

        # separator line
        _rect(s, Inches(4.0), Inches(0.95), Inches(0.02), Inches(0.8),
              RGBColor(0x33, 0x4A, 0x70))

        # 期初 / KO / KI (right block, 3 columns)
        col_x = [4.3, 7.0, 9.8]
        labels = [
            ("● 期初", f"${float(init_p):,.2f}" if _is_valid(init_p) else "N/A", _C_INIT),
            ("● KO",
             f"${ko_price:,.2f}" if ko_price and _is_valid(init_p) and True else "N/A",
             _C_KO),
            ("● KI",
             f"${ki_price:,.2f}" if ki_price and _is_valid(init_p) and True else "N/A",
             _C_KI),
        ]
        for (lbl, val, col), x in zip(labels, col_x):
            _textbox(s, Inches(x), Inches(0.86), Inches(2.8), Inches(0.38),
                     lbl, 12, bold=True, color=col)
            _textbox(s, Inches(x + 0.05), Inches(1.20), Inches(2.8), Inches(0.55),
                     val, 18, bold=True,
                     color=_WHITE if val == "N/A" else col)

        # ── second info row: 執行價 / 配息 / 比價日 / 出場日 ────────
        strike      = info.get("strike_pct")
        coupon      = info.get("coupon_pct")
        obs_date    = str(info.get("observation_date") or "")[:10]
        exit_date_v = str(info.get("exit_date") or "")[:10]

        customer_name = info.get("customer_name", "")
        amount_usd    = info.get("amount_usd")

        row2_items = []
        if customer_name:
            row2_items.append(("投資人", str(customer_name)))
        if _is_valid(amount_usd):
            row2_items.append(("金額", f"USD {float(amount_usd):,.0f}"))
        if _is_valid(strike):
            row2_items.append(("執行價", f"{float(strike)*100:.1f}%"))
        if _is_valid(coupon):
            row2_items.append(("配息", f"{float(coupon)*100:.2f}%"))
        if obs_date and obs_date != "None":
            row2_items.append(("比價日", obs_date))
        if exit_date_v and exit_date_v != "None":
            row2_items.append(("出場日", exit_date_v))
        if code and str(code).strip():
            row2_items.append(("商品", str(code)))

        if row2_items:
            _rect(s, 0, Inches(1.82), W, Inches(0.42), RGBColor(0x1A, 0x33, 0x5C))
            x_step = 13.33 / max(len(row2_items), 1)
            for idx, (lbl, val) in enumerate(row2_items):
                x = idx * x_step + 0.2
                _textbox(s, Inches(x), Inches(1.84), Inches(x_step - 0.1), Inches(0.18),
                         lbl, 8, color=RGBColor(0x94, 0xA3, 0xB8))
                _textbox(s, Inches(x), Inches(1.99), Inches(x_step - 0.1), Inches(0.22),
                         val, 10, bold=True, color=_WHITE)

        chart_top = Inches(2.28) if row2_items else Inches(1.95)
        chart_h   = Inches(5.1)  if row2_items else Inches(5.42)

        # ── hlines for chart ──────────────────────────────────────
        hlines = {}
        if _is_valid(init_p):
            hlines["initial"] = float(init_p)
            if ko_price:
                hlines["ko"] = ko_price
            if ki_price:
                hlines["ki"] = ki_price

        img_bytes = _make_chart_png(ticker, period=period, hlines=hlines or None)
        if img_bytes:
            img_stream = io.BytesIO(img_bytes)
            s.shapes.add_picture(img_stream,
                                 Inches(0.1), chart_top,
                                 Inches(13.1), chart_h)
        else:
            _textbox(s, Inches(1), Inches(4), Inches(11), Inches(1),
                     f"無法取得 {ticker} 圖表資料",
                     18, color=_GRAY, align=PP_ALIGN.CENTER)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()
